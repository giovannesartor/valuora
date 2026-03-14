"""
Quanto Vale — Pitch Deck PPTX Generator
Exports a pitch deck as a PowerPoint (.pptx) file using python-pptx.
Mirrors the PDF sections: Cover, Problem/Solution, Business Model, Market,
Team, Financial Plan, Funding.
"""
import io
import logging
from typing import Any, Dict, Optional

logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("[PPTX] python-pptx not installed — PPTX export unavailable.")


# ─── Brand Colors ────────────────────────────────────────
NAVY = RGBColor(0x0F, 0x17, 0x2A) if PPTX_AVAILABLE else None
EMERALD = RGBColor(0x10, 0xB9, 0x81) if PPTX_AVAILABLE else None
WHITE = RGBColor(0xFF, 0xFF, 0xFF) if PPTX_AVAILABLE else None
GRAY = RGBColor(0x4B, 0x55, 0x63) if PPTX_AVAILABLE else None
LIGHT_GRAY = RGBColor(0xE5, 0xE7, 0xEB) if PPTX_AVAILABLE else None


def _fmt_brl(value: Optional[float]) -> str:
    if value is None:
        return "—"
    if abs(value) >= 1_000_000:
        return f"R$ {value / 1_000_000:,.1f}M"
    elif abs(value) >= 1_000:
        return f"R$ {value / 1_000:,.0f}K"
    return f"R$ {value:,.0f}"


def _add_slide(prs: "Presentation", layout_idx: int = 6) -> Any:
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _fill_bg(slide: Any, color: "RGBColor") -> None:
    """Set slide background color."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree  # type: ignore
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _text_box(slide: Any, text: str, left: float, top: float, width: float, height: float,
              font_size: int = 18, bold: bool = False, color: "RGBColor" = None,
              align: str = "left") -> Any:
    """Add a text box to the slide."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return txBox


def _add_rect(slide: Any, left: float, top: float, width: float, height: float,
              fill_color: "RGBColor", line_color: "RGBColor" = None) -> Any:
    """Add a colored rectangle."""
    from pptx.util import Inches
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def _slide_cover(prs: "Presentation", deck: Any, analysis_data: Dict) -> None:
    slide = _add_slide(prs)
    _fill_bg(slide, NAVY)

    # Green accent bar on left
    _add_rect(slide, 0, 0, 0.08, 7.5, EMERALD)

    # Company name
    _text_box(slide, deck.company_name or "Empresa", 0.3, 0.6, 8.5, 1.2,
              font_size=40, bold=True, color=WHITE)

    # Slogan / AI headline
    headline = deck.ai_headline or deck.slogan or ""
    if headline:
        _text_box(slide, headline, 0.3, 1.9, 8.5, 0.8, font_size=16, color=EMERALD)

    # Sector + contact
    meta = f"{deck.sector or ''}"
    if deck.contact_email:
        meta += f"  ·  {deck.contact_email}"
    _text_box(slide, meta, 0.3, 6.5, 8, 0.6, font_size=11, color=GRAY)

    # Equity value if available
    ev = analysis_data.get("equity_value")
    if ev:
        _text_box(slide, "Valuation", 0.3, 3.0, 3, 0.4, font_size=10, color=EMERALD)
        _text_box(slide, _fmt_brl(ev), 0.3, 3.4, 4, 0.9, font_size=28, bold=True, color=WHITE)

    # Powered by badge
    _text_box(slide, "Valued by Valuora", 0.3, 7.0, 5, 0.4, font_size=9, color=GRAY)


def _slide_problem_solution(prs: "Presentation", deck: Any) -> None:
    slide = _add_slide(prs)
    _fill_bg(slide, WHITE)

    _add_rect(slide, 0, 0, 10, 0.08, EMERALD)
    _text_box(slide, "Problem & Solution", 0.4, 0.25, 9, 0.6, font_size=22, bold=True, color=NAVY)

    # Problem (left column)
    _text_box(slide, "Problem", 0.4, 1.0, 4.2, 0.4, font_size=13, bold=True, color=EMERALD)
    problem_text = deck.ai_problem or deck.problem or "Not defined."
    _text_box(slide, problem_text[:500], 0.4, 1.5, 4.2, 4.5, font_size=11, color=GRAY)

    # Solution (right column)
    _text_box(slide, "Solution", 5.2, 1.0, 4.2, 0.4, font_size=13, bold=True, color=EMERALD)
    solution_text = deck.ai_solution or deck.solution or "Not defined."
    _text_box(slide, solution_text[:500], 5.2, 1.5, 4.2, 4.5, font_size=11, color=GRAY)


def _slide_business_model(prs: "Presentation", deck: Any) -> None:
    slide = _add_slide(prs)
    _fill_bg(slide, WHITE)

    _add_rect(slide, 0, 0, 10, 0.08, EMERALD)
    _text_box(slide, "Business Model", 0.4, 0.25, 9, 0.6, font_size=22, bold=True, color=NAVY)

    bm = deck.ai_business_model or deck.business_model or "Not defined."
    _text_box(slide, bm[:800], 0.4, 1.1, 9.2, 5.5, font_size=12, color=GRAY)


def _slide_market(prs: "Presentation", deck: Any) -> None:
    slide = _add_slide(prs)
    _fill_bg(slide, NAVY)

    _add_rect(slide, 0, 0, 10, 0.08, EMERALD)
    _text_box(slide, "Market & Opportunity", 0.4, 0.25, 9, 0.6, font_size=22, bold=True, color=WHITE)

    tm = deck.target_market or {}
    if isinstance(tm, dict):
        desc = tm.get("description", "")
        if desc:
            _text_box(slide, desc[:400], 0.4, 1.1, 9.2, 2.5, font_size=13, color=LIGHT_GRAY)

        # TAM/SAM/SOM
        metrics = [
            ("TAM", tm.get("tam", "—")),
            ("SAM", tm.get("sam", "—")),
            ("SOM", tm.get("som", "—")),
        ]
        for i, (label, val) in enumerate(metrics):
            x = 0.4 + i * 3.2
            _add_rect(slide, x, 3.8, 2.8, 1.5, RGBColor(0x1E, 0x29, 0x3B))
            _text_box(slide, label, x + 0.1, 3.9, 2.6, 0.4, font_size=11, bold=True, color=EMERALD)
            _text_box(slide, str(val)[:30], x + 0.1, 4.35, 2.6, 0.7, font_size=13, color=WHITE)


def _slide_financial(prs: "Presentation", deck: Any, analysis_data: Dict) -> None:
    slide = _add_slide(prs)
    _fill_bg(slide, WHITE)

    _add_rect(slide, 0, 0, 10, 0.08, EMERALD)
    _text_box(slide, "Financial Plan", 0.4, 0.25, 9, 0.6, font_size=22, bold=True, color=NAVY)

    metrics = [
        ("Revenue (last year)", _fmt_brl(analysis_data.get("revenue"))),
        ("Net Margin", f"{(analysis_data.get('net_margin') or 0) * 100:.1f}%"),
        ("EBITDA", _fmt_brl(analysis_data.get("ebitda"))),
        ("Growth (projected)", f"{(analysis_data.get('growth_rate') or 0) * 100:.1f}%/ano"),
        ("Equity Value (Valuora)", _fmt_brl(analysis_data.get("equity_value"))),
        ("Risk Score", f"{analysis_data.get('risk_score') or '—'}/100"),
    ]

    for i, (label, val) in enumerate(metrics):
        row = i // 3
        col = i % 3
        x = 0.3 + col * 3.3
        y = 1.3 + row * 2.2
        _add_rect(slide, x, y, 3.0, 1.8, RGBColor(0xF9, 0xFA, 0xFB))
        _text_box(slide, label, x + 0.1, y + 0.1, 2.8, 0.5, font_size=9, color=GRAY)
        _text_box(slide, val, x + 0.1, y + 0.6, 2.8, 0.75, font_size=17, bold=True, color=NAVY)


def _slide_funding(prs: "Presentation", deck: Any) -> None:
    slide = _add_slide(prs)
    _fill_bg(slide, NAVY)

    _add_rect(slide, 0, 0, 10, 0.08, EMERALD)
    _text_box(slide, "Fundraising & Use of Funds", 0.4, 0.25, 9, 0.6, font_size=22, bold=True, color=WHITE)

    fn = deck.funding_needs or {}
    if isinstance(fn, dict):
        amount = fn.get("amount", 0)
        if amount:
            _text_box(slide, "Funding Round", 0.4, 1.1, 5, 0.4, font_size=12, color=EMERALD)
            _text_box(slide, _fmt_brl(amount), 0.4, 1.55, 5, 0.8, font_size=30, bold=True, color=WHITE)

    funding_use = deck.ai_funding_use or (fn.get("description") if isinstance(fn, dict) else None) or ""
    if funding_use:
        _text_box(slide, funding_use[:600], 0.4, 2.8, 9.2, 4.0, font_size=11, color=LIGHT_GRAY)


def _slide_team(prs: "Presentation", deck: Any) -> None:
    team = deck.team or []
    if not team:
        return

    slide = _add_slide(prs)
    _fill_bg(slide, WHITE)
    _add_rect(slide, 0, 0, 10, 0.08, EMERALD)
    _text_box(slide, "Founding Team", 0.4, 0.25, 9, 0.6, font_size=22, bold=True, color=NAVY)

    for i, member in enumerate(team[:4]):
        if isinstance(member, dict):
            col = i % 2
            row = i // 2
            x = 0.4 + col * 4.8
            y = 1.2 + row * 2.8
            _add_rect(slide, x, y, 4.3, 2.4, RGBColor(0xF3, 0xF4, 0xF6))
            _text_box(slide, member.get("name", ""), x + 0.15, y + 0.1, 4.0, 0.5, font_size=14, bold=True, color=NAVY)
            _text_box(slide, member.get("role", ""), x + 0.15, y + 0.65, 4.0, 0.4, font_size=10, bold=True, color=EMERALD)
            bio = member.get("bio", "")[:150] if member.get("bio") else ""
            if bio:
                _text_box(slide, bio, x + 0.15, y + 1.1, 4.0, 1.2, font_size=9, color=GRAY)


def generate_pitch_deck_pptx(deck: Any, analysis_data: Dict) -> bytes:
    """Generate a complete pitch deck as a PPTX file.

    Args:
        deck: PitchDeck ORM model instance
        analysis_data: Dict with financial data from linked analysis

    Returns:
        bytes: PPTX file contents
    """
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed. Run: pip install python-pptx")

    prs = Presentation()
    # Use widescreen 16:9 format
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Add slides
    _slide_cover(prs, deck, analysis_data)
    _slide_problem_solution(prs, deck)
    _slide_business_model(prs, deck)
    _slide_market(prs, deck)

    if deck.team:
        _slide_team(prs, deck)

    _slide_financial(prs, deck, analysis_data)
    _slide_funding(prs, deck)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

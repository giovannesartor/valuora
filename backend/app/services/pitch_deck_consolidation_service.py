"""
Pitch Deck Consolidation Service
Generates an executive consolidated report from N invites/pitch decks.

Pipeline:
  1) For each invite, extract structured summary via DeepSeek (with hash cache).
  2) Run cross-analysis (rankings, synergies, top picks).
  3) Render executive PDF (cover + cross-summary + per-deck sheet).
  4) Optionally render comparative PPTX.
  5) Optionally send via email (individual sends for privacy).
"""
from __future__ import annotations

import io
import os
import json
import re
import uuid
import hashlib
import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Optional, List, Dict

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import settings
from app.models.models import (
    PitchDeckInvite, PitchDeckConsolidation, PitchDeckConsolidationStatus,
)
from app.services.deepseek_service import call_deepseek
from app.services.pitch_deck_invite_service import compute_invite_score
from app.services.email_service import send_email

logger = logging.getLogger(__name__)

MAX_DECKS_PER_CONSOLIDATION = 10
SUMMARY_FIELDS = (
    "thesis", "problem_solution", "market", "traction", "team",
    "business_model", "competitive_edge", "ask",
    "risks", "recommendation",
)


def _submission_hash(submission: Optional[dict]) -> str:
    raw = json.dumps(submission or {}, sort_keys=True, ensure_ascii=False, default=str)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:32]


async def suggest_groups(invites: List[PitchDeckInvite]) -> dict:
    """Suggests groupings by sector/thesis/stage. Returns dict with GroupSuggestionResponse schema."""
    if not invites or len(invites) < 2:
        return {"groups": [], "ungrouped": [str(i.id) for i in invites]}

    compact = []
    for inv in invites:
        sub = inv.submission_data if isinstance(inv.submission_data, dict) else {}
        compact.append({
            "id": str(inv.id),
            "company": (inv.company_hint or sub.get("company_name") or sub.get("company") or "—")[:80],
            "sector": (sub.get("sector") or sub.get("industry") or "—")[:60],
            "stage": (sub.get("stage") or sub.get("round") or "—")[:40],
            "thesis": (sub.get("thesis") or sub.get("one_liner") or sub.get("description") or "")[:240],
        })

    prompt = (
        "You are a VC partner organizing a pipeline. Group the pitches below by affinity "
        "(sector, thesis, stage or model) to facilitate comparable consolidated reports. "
        "Each group must have 2 or more pitches; unmatched pitches go in 'ungrouped'.\n"
        "Respond ONLY strict JSON, no ```:\n"
        "{\n"
        '  "groups": [\n'
        '     {"label":"Fintech B2B","rationale":"all SaaS for banks","invite_ids":["uuid1","uuid2"]}\n'
        "  ],\n"
        '  "ungrouped": ["uuid3"]\n'
        "}\n\n"
        f"--- PITCHES ({len(compact)}) ---\n{json.dumps(compact, ensure_ascii=False)}"
    )
    try:
        raw = await call_deepseek(prompt, max_tokens=1200)
    except Exception as exc:
        logger.error(f"[consolidation] suggest_groups failed: {exc}")
        return {"groups": [], "ungrouped": [str(i.id) for i in invites], "_ai_error": str(exc)[:300]}

    match = re.search(r"\{[\s\S]*\}", raw or "")
    if not match:
        return {"groups": [], "ungrouped": [str(i.id) for i in invites]}
    try:
        data = json.loads(match.group(0))
    except json.JSONDecodeError:
        return {"groups": [], "ungrouped": [str(i.id) for i in invites]}

    valid_ids = {str(i.id) for i in invites}
    groups_clean = []
    used: set = set()
    for g in (data.get("groups") or []):
        ids = [x for x in (g.get("invite_ids") or []) if x in valid_ids and x not in used]
        if len(ids) < 2:
            continue
        used.update(ids)
        groups_clean.append({
            "label": str(g.get("label") or "Group")[:80],
            "rationale": str(g.get("rationale") or "")[:240],
            "invite_ids": ids,
        })
    ungrouped = [x for x in valid_ids if x not in used]
    return {"groups": groups_clean, "ungrouped": ungrouped}


async def summarize_invite(invite: PitchDeckInvite, *, force: bool = False) -> dict:
    """Generates (or reuses cache) structured summary of 1 invite."""
    submission = invite.submission_data if isinstance(invite.submission_data, dict) else {}
    cur_hash = _submission_hash(submission)

    cached = invite.ai_summary_json if isinstance(invite.ai_summary_json, dict) else None
    if cached and invite.ai_summary_hash == cur_hash and not force:
        cached["score"] = compute_invite_score(submission)
        return cached

    text = json.dumps(submission, ensure_ascii=False, default=str)[:14000]
    prompt = (
        "You are a senior venture capital analyst. Summarize the pitch below in strict JSON. "
        "Be objective and direct. If a field is missing, use empty string or empty list. "
        "DO NOT invent data — use only what appears in the pitch.\n\n"
        "Schema (respond ONLY with the JSON, no ```):\n"
        "{\n"
        '  "thesis": "1 sentence on the investment thesis",\n'
        '  "problem_solution": "max 2 lines: problem + solution",\n'
        '  "market": "TAM/SAM if available, short text",\n'
        '  "traction": "MRR/ARR/users/customers — real numbers only",\n'
        '  "team": "founders + relevant experience (max 2 lines)",\n'
        '  "business_model": "how it makes money (1 line)",\n'
        '  "competitive_edge": "defensible differentiator (1 line)",\n'
        '  "ask": "round + use of funds (1 line)",\n'
        '  "risks": ["risk 1", "risk 2", "risk 3"],\n'
        '  "recommendation": "pass | maybe | pass forward",\n'
        '  "recommendation_reason": "1 sentence justifying"\n'
        "}\n\n"
        f"--- PITCH ---\n{text}"
    )
    try:
        raw = await call_deepseek(prompt, max_tokens=900)
    except Exception as exc:
        logger.error(f"[consolidation] AI summarize failed for invite {invite.id}: {exc}")
        return {
            "thesis": "", "problem_solution": "", "market": "", "traction": "",
            "team": "", "business_model": "", "competitive_edge": "", "ask": "",
            "risks": [], "recommendation": "maybe",
            "recommendation_reason": "AI summary unavailable.",
            "score": compute_invite_score(submission),
            "_ai_error": str(exc)[:300],
        }

    match = re.search(r"\{[\s\S]*\}", raw or "")
    parsed: dict = {}
    if match:
        try:
            parsed = json.loads(match.group(0))
        except json.JSONDecodeError:
            parsed = {}
    if not parsed:
        parsed = {k: ("" if k != "risks" else []) for k in SUMMARY_FIELDS}

    for k in SUMMARY_FIELDS:
        parsed.setdefault(k, [] if k == "risks" else "")
    if not isinstance(parsed.get("risks"), list):
        parsed["risks"] = [str(parsed["risks"])] if parsed.get("risks") else []
    parsed["risks"] = [str(r).strip() for r in parsed["risks"] if r][:3]
    parsed["score"] = compute_invite_score(submission)

    invite.ai_summary_json = parsed
    invite.ai_summary_hash = cur_hash
    invite.ai_summary_generated_at = datetime.now(timezone.utc)
    return parsed


async def compute_meta_analysis(invites_with_summary: List[Dict[str, Any]]) -> dict:
    """Runs 1 AI call crossing all summaries."""
    if not invites_with_summary:
        return {}

    by_score = sorted(
        invites_with_summary,
        key=lambda x: int(((x.get("summary") or {}).get("score") or {}).get("score") or 0),
        reverse=True,
    )
    score_ranking = [
        {
            "id": x["id"],
            "company": x["company"],
            "score": int(((x.get("summary") or {}).get("score") or {}).get("score") or 0),
        }
        for x in by_score
    ]

    compact = [
        {
            "id": x["id"],
            "company": x["company"],
            "sector": x.get("sector") or "—",
            "thesis": x["summary"].get("thesis", ""),
            "traction": x["summary"].get("traction", ""),
            "competitive_edge": x["summary"].get("competitive_edge", ""),
            "ask": x["summary"].get("ask", ""),
            "score": int(((x.get("summary") or {}).get("score") or {}).get("score") or 0),
        }
        for x in invites_with_summary
    ]

    prompt = (
        "You are a venture capital partner analyzing a portfolio of opportunities. "
        "Compare the pitches below and produce strict JSON. Be objective and concise.\n\n"
        "Schema (respond ONLY with the JSON):\n"
        "{\n"
        '  "fit_ranking": [ {"id":"...", "company":"...", "fit":"high|medium|low", "reason":"1 sentence"} ],\n'
        '  "sector_overlap": "1-2 lines — is there sector overlap? which ones compete?",\n'
        '  "synergies": ["synergy 1 between A and B", "synergy 2..."],\n'
        '  "top_picks": [\n'
        '     {"id":"...", "company":"...", "reason":"why it is in the top 3"}\n'
        "  ],\n"
        '  "executive_takeaway": "1 paragraph (max 4 lines) with overall portfolio reading"\n'
        "}\n\n"
        f"--- PITCHES ({len(compact)}) ---\n{json.dumps(compact, ensure_ascii=False)}"
    )

    try:
        raw = await call_deepseek(prompt, max_tokens=1500)
    except Exception as exc:
        logger.error(f"[consolidation] meta-analysis failed: {exc}")
        return {
            "score_ranking": score_ranking,
            "fit_ranking": [],
            "sector_overlap": "",
            "synergies": [],
            "top_picks": score_ranking[:3],
            "executive_takeaway": "AI cross-analysis unavailable.",
            "_ai_error": str(exc)[:300],
        }

    match = re.search(r"\{[\s\S]*\}", raw or "")
    meta: dict = {}
    if match:
        try:
            meta = json.loads(match.group(0))
        except json.JSONDecodeError:
            meta = {}

    meta.setdefault("fit_ranking", [])
    meta.setdefault("sector_overlap", "")
    meta.setdefault("synergies", [])
    meta.setdefault("top_picks", [])
    meta.setdefault("executive_takeaway", "")
    meta["score_ranking"] = score_ranking
    return meta


def _ensure_dir() -> Path:
    out_dir = Path(getattr(settings, "REPORTS_DIR", "storage/reports"))
    out_dir.mkdir(parents=True, exist_ok=True)
    return out_dir


def render_consolidation_pdf(
    *,
    consolidation: PitchDeckConsolidation,
    invites: List[PitchDeckInvite],
    summaries: List[dict],
    meta: dict,
) -> str:
    """Generates the executive PDF. Returns the absolute path."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_LEFT, TA_CENTER
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, HRFlowable,
    )

    out_dir = _ensure_dir()
    fname = f"consolidation-{consolidation.id}.pdf"
    fpath = out_dir / fname

    NAVY = HexColor("#0f172a")
    EMERALD = HexColor("#10b981")
    GRAY_700 = HexColor("#374151")
    GRAY_500 = HexColor("#6b7280")
    GRAY_200 = HexColor("#e5e7eb")
    GRAY_50 = HexColor("#f9fafb")

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle("h1", parent=styles["Heading1"], fontSize=22, leading=26, textColor=NAVY, spaceAfter=10)
    h2 = ParagraphStyle("h2", parent=styles["Heading2"], fontSize=14, leading=18, textColor=EMERALD, spaceAfter=8, spaceBefore=10)
    h3 = ParagraphStyle("h3", parent=styles["Heading3"], fontSize=11, leading=14, textColor=NAVY, spaceAfter=4, spaceBefore=4)
    body = ParagraphStyle("body", parent=styles["BodyText"], fontSize=9.5, leading=13, textColor=GRAY_700)
    small = ParagraphStyle("small", parent=styles["BodyText"], fontSize=8, leading=11, textColor=GRAY_500)
    cover_title = ParagraphStyle("cover_title", parent=styles["Heading1"], fontSize=32, leading=38, textColor=NAVY, alignment=TA_CENTER)
    cover_sub = ParagraphStyle("cover_sub", parent=styles["BodyText"], fontSize=14, leading=18, textColor=GRAY_500, alignment=TA_CENTER)

    def _esc(s: Any) -> str:
        s = "" if s is None else str(s)
        return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

    story: list = []

    title = consolidation.title or f"Executive Consolidated Report — {len(invites)} pitches"
    story.append(Spacer(1, 5 * cm))
    story.append(Paragraph(_esc(title), cover_title))
    story.append(Spacer(1, 0.8 * cm))
    story.append(Paragraph(
        f"Comparative analysis of {len(invites)} opportunities · "
        f"Generated on {datetime.now().strftime('%m/%d/%Y %H:%M')}",
        cover_sub,
    ))
    story.append(Spacer(1, 4 * cm))
    story.append(Paragraph("Valuora · Smart Investor", cover_sub))
    story.append(PageBreak())

    story.append(Paragraph("Executive Summary", h1))
    story.append(HRFlowable(width="100%", thickness=1, color=EMERALD, spaceAfter=10))
    if meta.get("executive_takeaway"):
        story.append(Paragraph(_esc(meta["executive_takeaway"]), body))
        story.append(Spacer(1, 0.3 * cm))

    story.append(Paragraph("Score Ranking", h2))
    rows = [["#", "Company", "Score"]]
    for i, item in enumerate(meta.get("score_ranking") or [], 1):
        rows.append([str(i), _esc(item.get("company")), str(item.get("score", 0))])
    if len(rows) > 1:
        t = Table(rows, colWidths=[1 * cm, 12 * cm, 2.5 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), NAVY),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#ffffff")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [GRAY_50, HexColor("#ffffff")]),
            ("GRID", (0, 0), (-1, -1), 0.3, GRAY_200),
            ("ALIGN", (0, 0), (0, -1), "CENTER"),
            ("ALIGN", (2, 0), (2, -1), "CENTER"),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(t)
        story.append(Spacer(1, 0.4 * cm))

    if meta.get("fit_ranking"):
        story.append(Paragraph("Investment Fit Ranking (AI)", h2))
        rows = [["Company", "Fit", "Reason"]]
        for item in meta["fit_ranking"]:
            rows.append([
                _esc(item.get("company")),
                _esc(item.get("fit")),
                Paragraph(_esc(item.get("reason")), body),
            ])
        t = Table(rows, colWidths=[5 * cm, 2 * cm, 8.5 * cm])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), NAVY),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#ffffff")),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [GRAY_50, HexColor("#ffffff")]),
            ("GRID", (0, 0), (-1, -1), 0.3, GRAY_200),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(t)
        story.append(Spacer(1, 0.4 * cm))

    if meta.get("sector_overlap"):
        story.append(Paragraph("Sector Comparison", h2))
        story.append(Paragraph(_esc(meta["sector_overlap"]), body))
        story.append(Spacer(1, 0.3 * cm))

    if meta.get("synergies"):
        story.append(Paragraph("Possible Synergies", h2))
        for s in meta["synergies"]:
            story.append(Paragraph(f"• {_esc(s)}", body))
        story.append(Spacer(1, 0.3 * cm))

    if meta.get("top_picks"):
        story.append(Paragraph("Top Picks", h2))
        h3_style = ParagraphStyle("h3", parent=styles["Heading3"], fontSize=11, leading=14, textColor=NAVY, spaceAfter=4, spaceBefore=4)
        for tp in meta["top_picks"][:3]:
            story.append(Paragraph(f"<b>{_esc(tp.get('company'))}</b>", h3_style))
            story.append(Paragraph(_esc(tp.get("reason") or ""), body))
            story.append(Spacer(1, 0.2 * cm))

    story.append(PageBreak())

    invites_by_id = {str(i.id): i for i in invites}
    for idx, summary in enumerate(summaries, 1):
        invite = invites_by_id.get(summary.get("_invite_id"))
        if not invite:
            continue
        sub = invite.submission_data if isinstance(invite.submission_data, dict) else {}
        company = sub.get("company_name") or invite.company_hint or invite.client_name or "—"
        sector = sub.get("sector") or "—"
        score_obj = summary.get("score") or {}
        score_n = score_obj.get("score", 0)
        rec = (summary.get("recommendation") or "—").upper()

        story.append(Paragraph(f"{idx}. {_esc(company)}", h1))
        story.append(Paragraph(_esc(f"Sector: {sector} · Score: {score_n}/100 · Recommendation: {rec}"), small))
        story.append(HRFlowable(width="100%", thickness=0.6, color=GRAY_200, spaceAfter=8))

        h3_s = ParagraphStyle("h3s", parent=styles["Heading3"], fontSize=11, leading=14, textColor=NAVY, spaceAfter=4, spaceBefore=4)

        def _section(label: str, value: str):
            if not value:
                return
            story.append(Paragraph(label, h3_s))
            story.append(Paragraph(_esc(value), body))
            story.append(Spacer(1, 0.15 * cm))

        _section("Thesis", summary.get("thesis"))
        _section("Problem & Solution", summary.get("problem_solution"))
        _section("Market", summary.get("market"))
        _section("Traction", summary.get("traction"))
        _section("Team", summary.get("team"))
        _section("Business Model", summary.get("business_model"))
        _section("Competitive Edge", summary.get("competitive_edge"))
        _section("Ask (round)", summary.get("ask"))

        risks = summary.get("risks") or []
        if risks:
            story.append(Paragraph("Top 3 Risks", h3_s))
            for r in risks[:3]:
                story.append(Paragraph(f"• {_esc(r)}", body))
            story.append(Spacer(1, 0.15 * cm))

        if summary.get("recommendation_reason"):
            story.append(Paragraph("Recommendation Justification", h3_s))
            story.append(Paragraph(_esc(summary["recommendation_reason"]), body))

        story.append(PageBreak())

    doc = SimpleDocTemplate(
        str(fpath),
        pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm,
        topMargin=2 * cm, bottomMargin=2 * cm,
        title=title, author="Valuora", subject="Pitch Deck Consolidation",
    )
    doc.build(story)
    return str(fpath)


def render_consolidation_pptx(
    *,
    consolidation: PitchDeckConsolidation,
    invites: List[PitchDeckInvite],
    summaries: List[dict],
    meta: dict,
) -> Optional[str]:
    """Generates comparative PPTX. Returns path or None if python-pptx unavailable."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.warning("[consolidation] python-pptx not installed — skipping PPTX")
        return None

    out_dir = _ensure_dir()
    fname = f"consolidation-{consolidation.id}.pptx"
    fpath = out_dir / fname

    NAVY = RGBColor(0x0F, 0x17, 0x2A)
    EMERALD = RGBColor(0x10, 0xB9, 0x81)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x4B, 0x55, 0x63)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def _bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _txt(slide, text, left, top, width, height, *, size=18, bold=False, color=NAVY, align="left"):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        return box

    blank = prs.slide_layouts[6]

    # Cover
    s = prs.slides.add_slide(blank)
    _bg(s, NAVY)
    _txt(s, consolidation.title or f"Consolidated Report — {len(invites)} pitches",
         0.5, 2.5, 12.3, 1.2, size=36, bold=True, color=WHITE, align="center")
    _txt(s, f"Comparative analysis · {datetime.now().strftime('%m/%d/%Y')}",
         0.5, 4.0, 12.3, 0.6, size=18, color=EMERALD, align="center")
    _txt(s, "Valuora", 0.5, 6.5, 12.3, 0.5, size=12, color=WHITE, align="center")

    # Ranking slide
    s = prs.slides.add_slide(blank)
    _bg(s, WHITE)
    _txt(s, "Score Ranking", 0.6, 0.4, 12, 0.7, size=28, bold=True, color=NAVY)
    y = 1.3
    for i, item in enumerate((meta.get("score_ranking") or [])[:10], 1):
        _txt(s, f"{i}. {item.get('company')}", 0.8, y, 8, 0.4, size=16, bold=True, color=NAVY)
        _txt(s, f"{item.get('score', 0)}/100", 9.5, y, 2, 0.4, size=16, color=EMERALD, align="center")
        y += 0.55

    if meta.get("top_picks"):
        s = prs.slides.add_slide(blank)
        _bg(s, WHITE)
        _txt(s, "Top Picks", 0.6, 0.4, 12, 0.7, size=28, bold=True, color=NAVY)
        y = 1.4
        for tp in meta["top_picks"][:3]:
            _txt(s, tp.get("company") or "", 0.8, y, 12, 0.5, size=20, bold=True, color=EMERALD)
            _txt(s, tp.get("reason") or "", 0.8, y + 0.5, 12, 1.2, size=14, color=GRAY)
            y += 1.9

    invites_by_id = {str(i.id): i for i in invites}
    for summary in summaries:
        invite = invites_by_id.get(summary.get("_invite_id"))
        if not invite:
            continue
        sub = invite.submission_data if isinstance(invite.submission_data, dict) else {}
        company = sub.get("company_name") or invite.company_hint or invite.client_name or "—"
        score_n = (summary.get("score") or {}).get("score", 0)

        s = prs.slides.add_slide(blank)
        _bg(s, WHITE)
        _txt(s, company, 0.6, 0.4, 9, 0.8, size=28, bold=True, color=NAVY)
        _txt(s, f"Score {score_n}/100", 10.0, 0.5, 3, 0.6, size=18, bold=True, color=EMERALD, align="center")
        _txt(s, sub.get("sector") or "", 0.6, 1.1, 9, 0.4, size=12, color=GRAY)

        rows = [
            ("Thesis", summary.get("thesis", "")),
            ("Problem & Solution", summary.get("problem_solution", "")),
            ("Market", summary.get("market", "")),
            ("Traction", summary.get("traction", "")),
            ("Team", summary.get("team", "")),
            ("Model", summary.get("business_model", "")),
            ("Edge", summary.get("competitive_edge", "")),
            ("Ask", summary.get("ask", "")),
        ]
        y = 1.8
        for label, val in rows:
            if not val:
                continue
            _txt(s, label, 0.6, y, 2.2, 0.35, size=11, bold=True, color=EMERALD)
            _txt(s, val, 2.9, y, 10.0, 0.45, size=11, color=NAVY)
            y += 0.55
            if y > 6.8:
                break

        rec = summary.get("recommendation") or "—"
        _txt(s, f"Recommendation: {rec.upper()}", 0.6, 7.0, 12, 0.4, size=12, bold=True, color=GRAY)

    prs.save(str(fpath))
    return str(fpath)


async def email_consolidation(
    *,
    consolidation: PitchDeckConsolidation,
    recipients: List[str],
) -> bool:
    """Sends PDF (and PPTX if available) by email to each recipient individually."""
    if not recipients or not consolidation.pdf_path:
        return False
    if not os.path.exists(consolidation.pdf_path):
        logger.error(f"[consolidation] PDF not found: {consolidation.pdf_path}")
        return False

    title = consolidation.title or "Consolidated Pitch Deck Report"
    html = (
        f"<p>Hello,</p>"
        f"<p>Please find attached the executive consolidated report: <b>{title}</b>.</p>"
        f"<p>Generated on {datetime.now().strftime('%m/%d/%Y %H:%M')} by Valuora.</p>"
    )

    sent_any = False
    for addr in recipients:
        try:
            await send_email(
                to_email=addr,
                subject=f"[Valuora] {title}",
                html_body=html,
                attachment_path=consolidation.pdf_path,
            )
            sent_any = True
        except Exception as exc:
            logger.warning(f"[consolidation] email to {addr} failed: {exc}")
    return sent_any


async def run_consolidation(
    db: AsyncSession,
    consolidation_id: uuid.UUID,
) -> None:
    """Main job: summaries → meta → PDF → PPTX (opt) → email (opt)."""
    res = await db.execute(select(PitchDeckConsolidation).where(PitchDeckConsolidation.id == consolidation_id))
    cons: Optional[PitchDeckConsolidation] = res.scalar_one_or_none()
    if not cons:
        logger.error(f"[consolidation] not found: {consolidation_id}")
        return
    if cons.deleted_at:
        return

    cons.status = PitchDeckConsolidationStatus.PROCESSING
    cons.progress_pct = 5
    cons.progress_message = "Loading pitches..."
    cons.error = None
    await db.commit()

    try:
        invite_ids = [uuid.UUID(x) for x in (cons.invite_ids or [])]
        if not invite_ids:
            raise ValueError("No pitch decks selected.")

        res = await db.execute(
            select(PitchDeckInvite).where(
                PitchDeckInvite.id.in_(invite_ids),
                PitchDeckInvite.deleted_at.is_(None),
            )
        )
        invites = res.scalars().all()
        if not invites:
            raise ValueError("No valid pitch decks found.")

        summaries: list[dict] = []
        n = len(invites)
        for i, inv in enumerate(invites):
            cons.progress_pct = 10 + int(60 * i / max(n, 1))
            cons.progress_message = f"Summarizing {i + 1}/{n}..."
            await db.commit()

            s = await summarize_invite(inv)
            s = dict(s)
            s["_invite_id"] = str(inv.id)
            summaries.append(s)

        await db.commit()

        cons.progress_pct = 75
        cons.progress_message = "Cross-analysis..."
        await db.commit()

        items_for_meta = []
        for inv, s in zip(invites, summaries):
            sub = inv.submission_data if isinstance(inv.submission_data, dict) else {}
            items_for_meta.append({
                "id": str(inv.id),
                "company": sub.get("company_name") or inv.company_hint or inv.client_name or "—",
                "sector": sub.get("sector") or "",
                "summary": s,
            })
        meta = await compute_meta_analysis(items_for_meta)

        cons.progress_pct = 85
        cons.progress_message = "Generating PDF..."
        await db.commit()
        pdf_path = render_consolidation_pdf(
            consolidation=cons, invites=list(invites), summaries=summaries, meta=meta,
        )
        cons.pdf_path = pdf_path
        cons.meta_json = meta

        opts = cons.options or {}
        if opts.get("include_pptx"):
            cons.progress_pct = 92
            cons.progress_message = "Generating PPTX..."
            await db.commit()
            try:
                pptx_path = render_consolidation_pptx(
                    consolidation=cons, invites=list(invites), summaries=summaries, meta=meta,
                )
                if pptx_path:
                    cons.pptx_path = pptx_path
            except Exception as exc:
                logger.warning(f"[consolidation] PPTX failed: {exc}")

        recipients = opts.get("email_recipients") or []
        if recipients:
            cons.progress_pct = 97
            cons.progress_message = "Sending email..."
            await db.commit()
            try:
                await email_consolidation(consolidation=cons, recipients=recipients)
            except Exception as exc:
                logger.warning(f"[consolidation] email failed: {exc}")

        cons.status = PitchDeckConsolidationStatus.READY
        cons.progress_pct = 100
        cons.progress_message = "Done."
        cons.ready_at = datetime.now(timezone.utc)
        await db.commit()
    except Exception as exc:
        logger.exception(f"[consolidation] failed: {exc}")
        cons.status = PitchDeckConsolidationStatus.FAILED
        cons.error = str(exc)[:1000]
        cons.progress_message = "Failed."
        await db.commit()
